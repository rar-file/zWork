"""Tool execution for the chat agent.

Each tool is:
  - declared in TOOL_SCHEMAS for the LLM
  - implemented here with a clear contract:
      yields `activity` events (for the UI)
      yields exactly one `tool_result` event at the end
"""
from __future__ import annotations

import json
import os
import re
import sys
import subprocess
import shlex
from pathlib import Path
from typing import Any, AsyncIterator


# ---------------- Tool schemas (provider-neutral) ----------------

from . import skills as skills_mod
from .home import memory_path


TOOL_SCHEMAS: list[dict] = [
    {
        "name": "write_file",
        "description": (
            "Write content to a file at the given path. Creates parent directories if needed. "
            "Overwrites existing files. Use this for creating app files, code, docs, configs."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Relative or absolute file path"},
                "content": {"type": "string", "description": "Full UTF-8 content of the file"},
            },
            "required": ["path", "content"],
        },
    },
    {
        "name": "read_file",
        "description": (
            "Read and return the UTF-8 contents of a file. Use this to inspect files you or "
            "the user just wrote/changed, or to see existing project files before editing."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Relative or absolute file path"},
            },
            "required": ["path"],
        },
    },
    {
        "name": "list_dir",
        "description": "List the immediate children of a directory. Use this to orient yourself.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Directory path (default: '.')"},
            },
            "required": [],
        },
    },
    {
        "name": "run_command",
        "description": (
            "Run a shell command. Set background=true for long-running servers (e.g. dev servers); "
            "the command will detach and return immediately. For foreground commands, the combined "
            "stdout+stderr is returned (120s timeout)."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "command": {"type": "string", "description": "Shell command to execute"},
                "cwd": {"type": "string", "description": "Working directory (default: '.')"},
                "background": {"type": "boolean", "description": "Run detached (for servers)"},
            },
            "required": ["command"],
        },
    },
    {
        "name": "read_skill",
        "description": (
            "Load the full SKILL.md for an installed skill so you can follow its playbook. "
            "Pass the skill slug (e.g. 'anthropic-skills/pdf' or 'uiux-pro-max'). "
            "The system prompt lists available skills. Call this when a user task matches a skill's domain."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "slug": {"type": "string", "description": "Skill slug or leaf folder name"},
            },
            "required": ["slug"],
        },
    },
    {
        "name": "deploy_web_app",
        "description": (
            "Serve a local web app directory on http://localhost:<port>. Picks a free port "
            "(prefers 5173, 8000, 3000). Uses `python3 -m http.server` for static sites, or "
            "`npm run dev` when package.json has a dev script. Returns the URL."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "project_path": {"type": "string", "description": "Path to the project root"},
                "framework": {"type": "string", "description": "Framework hint (optional)"},
            },
            "required": ["project_path"],
        },
    },
    {
        "name": "save_memory",
        "description": (
            "Append a fact or note to persistent memory. "
            "MUST be called whenever the user asks to remember, note down, or save something. "
            "Do NOT just acknowledge the request — actually invoke this tool."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "content": {"type": "string", "description": "The fact or note to remember"},
            },
            "required": ["content"],
        },
    },
    {
        "name": "extract_document",
        "description": (
            "Extract text (and where applicable, tables and metadata) from a document on disk. "
            "Supports PDF, DOCX, XLSX, PPTX, TXT and Markdown — auto-detected from the file "
            "extension. Use this before answering questions about the contents of a file the "
            "user has pointed at; do not try to read these formats with read_file."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the document"},
                "format": {
                    "type": "string",
                    "description": "Output style hint ('markdown' or 'text'); default 'markdown'",
                },
                "pages": {
                    "type": "string",
                    "description": "Optional 1-based page range for PDFs, e.g. '1-5' or '3'",
                },
            },
            "required": ["path"],
        },
    },
    {
        "name": "dctl",
        "description": (
            "Run the local dctl desktop-control CLI for window/app/browser automation, "
            "accessibility tree inspection, screenshots, and focus/click/type/scroll actions. "
            "Prefer this over raw shell for GUI work."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "subcommand": {"type": "string", "description": "dctl subcommand to run"},
                "args": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Arguments passed to the dctl subcommand",
                },
                "cwd": {"type": "string", "description": "Working directory for the command"},
            },
            "required": ["subcommand"],
        },
    },
]


# ---------------- Dispatcher ----------------

def _friendly_error(err: Exception, context: str = "") -> str:
    """Translate common exceptions into user-actionable messages."""
    msg = str(err)

    if isinstance(err, FileNotFoundError):
        return f"File not found: {msg}. Check the path and try again."
    if isinstance(err, NotADirectoryError):
        return f"Not a directory: {msg}. Specify a directory path instead."
    if isinstance(err, IsADirectoryError):
        return f"Is a directory, not a file: {msg}. Specify a file path instead."
    if isinstance(err, PermissionError):
        return f"Permission denied: {msg}. Check file permissions and try again."
    if isinstance(err, UnicodeDecodeError):
        return f"Cannot read as text: {msg}. The file may be binary or use an unsupported encoding."
    if isinstance(err, subprocess.TimeoutExpired):
        return f"Command timed out after {err.timeout}s. Try breaking the work into smaller steps."
    if isinstance(err, OSError):
        return f"System error: {msg}. {context}Check that the path and permissions are correct."

    return f"{msg}. {context}" if context else msg


async def execute_tool(tool_name: str, params: dict[str, Any]) -> AsyncIterator[dict]:
    tool_id = f"tool_{tool_name}_{id(params)}"

    if tool_name == "write_file":
        path = params.get("path", "")
        content = params.get("content", "")
        label = f"Write {_short_path(path)}"
        icon = _icon_for_path(path)
        yield {"type": "activity", "id": tool_id, "label": label, "icon": icon, "done": False}
        try:
            _write_file(path, content)
            yield {"type": "activity", "id": tool_id, "label": label, "icon": icon, "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": True,
                   "message": f"Wrote {len(content)} chars to {path}"}
        except Exception as e:
            yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": icon, "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": False, "message": _friendly_error(e, "Try a different path. ")}
        return

    if tool_name == "read_file":
        path = params.get("path", "")
        label = f"Read {_short_path(path)}"
        icon = _icon_for_path(path)
        yield {"type": "activity", "id": tool_id, "label": label, "icon": icon, "done": False}
        try:
            text = _read_file(path)
            yield {"type": "activity", "id": tool_id, "label": label, "icon": icon, "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": True, "message": text}
        except Exception as e:
            yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": icon, "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": False, "message": _friendly_error(e, "Try a different path. ")}
        return

    if tool_name == "list_dir":
        path = params.get("path", ".")
        label = f"List {_short_path(path)}"
        yield {"type": "activity", "id": tool_id, "label": label, "icon": "folder", "done": False}
        try:
            listing = _list_dir(path)
            yield {"type": "activity", "id": tool_id, "label": label, "icon": "folder", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": True, "message": listing}
        except Exception as e:
            yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": "folder", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": False, "message": _friendly_error(e, "Check that the path is a directory. ")}
        return

    if tool_name == "run_command":
        command = params.get("command", "")
        cwd = params.get("cwd", ".")
        background = bool(params.get("background", False))
        short = command[:60] + ("…" if len(command) > 60 else "")
        label = f"Run: {short}" + (" (bg)" if background else "")
        yield {"type": "activity", "id": tool_id, "label": label, "icon": "command", "done": False}
        try:
            if background:
                pid = _run_background(command, cwd)
                yield {"type": "activity", "id": tool_id, "label": label, "icon": "command", "done": True}
                yield {"type": "tool_result", "tool": tool_name, "ok": True,
                       "message": f"Started background process (pid={pid})"}
            else:
                result = _run_command(command, cwd)
                yield {"type": "activity", "id": tool_id, "label": label, "icon": "command", "done": True}
                yield {"type": "tool_result", "tool": tool_name, "ok": result["ok"],
                       "message": result["output"] or (f"exit {result['returncode']}")}
        except Exception as e:
            yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": "command", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": False, "message": _friendly_error(e)}
        return

    if tool_name == "read_skill":
        slug = params.get("slug", "")
        label = f"Read skill {slug}"
        yield {"type": "activity", "id": tool_id, "label": label, "icon": "file", "done": False}
        try:
            text = skills_mod.read_skill(slug)
            if text is None:
                available = ", ".join(s.slug for s in skills_mod.list_skills()[:10])
                msg = f"No skill named '{slug}'. Try one of: {available}"
                yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": "file", "done": True}
                yield {"type": "tool_result", "tool": tool_name, "ok": False, "message": msg}
            else:
                # Cap to keep context sane.
                if len(text) > 80_000:
                    text = text[:80_000] + "\n…[truncated]"
                yield {"type": "activity", "id": tool_id, "label": label, "icon": "file", "done": True}
                yield {"type": "tool_result", "tool": tool_name, "ok": True, "message": text}
        except Exception as e:
            yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": "file", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": False, "message": _friendly_error(e, "Check the skill slug spelling. ")}
        return

    if tool_name == "deploy_web_app":
        project_path = params.get("project_path", ".")
        framework = params.get("framework", "")
        label = f"Serve {_short_path(project_path)}"
        yield {"type": "activity", "id": tool_id, "label": label, "icon": "deploy", "done": False}
        try:
            result = _deploy_web_app(project_path, framework)
            yield {"type": "activity", "id": tool_id, "label": label, "icon": "deploy", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": result["ok"],
                   "message": result["message"]}
        except Exception as e:
            yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": "deploy", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": False, "message": _friendly_error(e)}
        return

    if tool_name == "save_memory":
        content = params.get("content", "")
        label = "Save memory"
        yield {"type": "activity", "id": tool_id, "label": label, "icon": "file", "done": False}
        try:
            _save_memory(content)
            yield {"type": "activity", "id": tool_id, "label": label, "icon": "file", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": True,
                   "message": "Saved to memory."}
        except Exception as e:
            yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": "file", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": False, "message": _friendly_error(e)}
        return

    if tool_name == "extract_document":
        path = params.get("path", "")
        fmt = params.get("format", "markdown") or "markdown"
        pages = params.get("pages")
        label = f"Extract {_short_path(path)}"
        yield {"type": "activity", "id": tool_id, "label": label, "icon": "file", "done": False}
        try:
            result = _extract_document(path, fmt, pages)
            yield {"type": "activity", "id": tool_id, "label": label, "icon": "file", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": True,
                   "message": json.dumps(result, ensure_ascii=False)}
        except Exception as e:
            yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": "file", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": False,
                   "message": _friendly_error(e, "Check the file path and extension. ")}
        return

    if tool_name == "dctl":
        subcommand = str(params.get("subcommand", "")).strip()
        args = [str(a) for a in (params.get("args") or []) if str(a).strip()]
        cwd = params.get("cwd", ".")
        full = " ".join([subcommand, *args]).strip()
        label = f"dctl {full}" if full else "dctl"
        yield {"type": "activity", "id": tool_id, "label": label, "icon": "window", "done": False}
        try:
            result = _run_dctl(subcommand, args, cwd)
            yield {"type": "activity", "id": tool_id, "label": label, "icon": "window", "done": True}
            yield {
                "type": "tool_result",
                "tool": tool_name,
                "ok": result["ok"],
                "message": result["output"] or (f"exit {result['returncode']}" if not result["ok"] else "ok"),
            }
        except Exception as e:
            yield {"type": "activity", "id": tool_id, "label": f"Failed: {label}", "icon": "window", "done": True}
            yield {"type": "tool_result", "tool": tool_name, "ok": False, "message": _friendly_error(e, "Check that dctl is installed and the subcommand is valid. ")}
        return

    yield {"type": "tool_result", "tool": tool_name, "ok": False,
           "message": f"Unknown tool: {tool_name}. This tool is not available — try a different approach."}


# ---------------- Impls ----------------

def _short_path(path: str) -> str:
    try:
        home = str(Path.home())
        s = str(path)
        if s.startswith(home):
            s = "~" + s[len(home):]
        return s
    except Exception:
        return str(path)


def _icon_for_path(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext in (".html", ".htm"):
        return "html"
    if ext in (".css",):
        return "css"
    if ext in (".js", ".mjs"):
        return "js"
    if ext in (".ts", ".tsx"):
        return "ts"
    if ext in (".jsx",):
        return "jsx"
    if ext in (".json",):
        return "json"
    if ext in (".py",):
        return "code"
    if ext in (".md", ".markdown"):
        return "file"
    return "file"


def _write_file(path: str, content: str) -> None:
    p = Path(path).expanduser()
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content, encoding="utf-8")


def _read_file(path: str) -> str:
    p = Path(path).expanduser()
    data = p.read_text(encoding="utf-8", errors="replace")
    # Cap to avoid flooding the context
    if len(data) > 200_000:
        return data[:200_000] + "\n…[truncated]"
    return data


def _list_dir(path: str) -> str:
    p = Path(path).expanduser()
    if not p.exists():
        raise FileNotFoundError(f"No such directory: {path}")
    if not p.is_dir():
        raise NotADirectoryError(f"Not a directory: {path}")
    entries = []
    for child in sorted(p.iterdir()):
        suffix = "/" if child.is_dir() else ""
        entries.append(child.name + suffix)
    return "\n".join(entries) if entries else "(empty)"


def _run_command(command: str, cwd: str) -> dict[str, Any]:
    try:
        result = subprocess.run(
            command,
            shell=True,
            cwd=Path(cwd).expanduser(),
            capture_output=True,
            text=True,
            timeout=120,
        )
    except subprocess.TimeoutExpired:
        return {
            "ok": False,
            "returncode": -1,
            "output": f"Command timed out after 120s. Try a shorter command or break the work into smaller steps.",
        }
    output = result.stdout
    if result.stderr:
        output += ("\n" + result.stderr) if output else result.stderr
    # Cap output
    if len(output) > 20_000:
        output = output[:20_000] + "\n…[truncated]"
    return {
        "ok": result.returncode == 0,
        "returncode": result.returncode,
        "output": output.strip(),
    }


def _run_background(command: str, cwd: str) -> int:
    """Start a detached background process. Returns PID."""
    proc = subprocess.Popen(
        command,
        shell=True,
        cwd=Path(cwd).expanduser(),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        stdin=subprocess.DEVNULL,
        start_new_session=True,
    )
    return proc.pid


def _pick_free_port(preferred: list[int]) -> int:
    import socket
    for port in preferred:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
            try:
                sock.bind(("127.0.0.1", port))
                return port
            except OSError:
                continue
    # fall back to kernel-assigned
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
        sock.bind(("127.0.0.1", 0))
        return sock.getsockname()[1]


def _deploy_web_app(project_path: str, framework: str) -> dict[str, Any]:
    """Actually start a server."""
    p = Path(project_path).expanduser().resolve()
    if not p.exists():
        return {"ok": False, "message": f"Project path does not exist: {project_path}"}
    if not p.is_dir():
        return {"ok": False, "message": f"Project path is not a directory: {project_path}"}

    pkg = p / "package.json"
    if pkg.exists():
        try:
            data = json.loads(pkg.read_text())
            scripts = data.get("scripts") or {}
            if "dev" in scripts:
                port = _pick_free_port([5173, 3000, 8080])
                # Pass PORT env for CRA/Next; Vite uses --port via the script itself if configured.
                command = f"PORT={port} npm run dev"
                _run_background(command, str(p))
                return {"ok": True,
                        "message": f"Started `npm run dev` in {p.name}. Open http://localhost:{port} "
                                   f"(check console if your dev server chose a different port)."}
            if "start" in scripts:
                port = _pick_free_port([3000, 8080])
                _run_background(f"PORT={port} npm start", str(p))
                return {"ok": True,
                        "message": f"Started `npm start` in {p.name}. Open http://localhost:{port}."}
        except Exception:
            pass

    index = p / "index.html"
    if index.exists():
        port = _pick_free_port([8000, 8080, 5173])
        _run_background(f"python3 -m http.server {port}", str(p))
        return {"ok": True,
                "message": f"Serving {p.name} at http://localhost:{port}"}

    return {"ok": False,
            "message": f"No index.html or package.json in {p}. Nothing obvious to serve."}


def _save_memory(content: str) -> None:
    """Append a note to the global memory file."""
    import time as _time
    p = memory_path()
    p.parent.mkdir(parents=True, exist_ok=True)
    existing = ""
    if p.exists():
        existing = p.read_text(encoding="utf-8").rstrip()
    timestamp = _time.strftime("%Y-%m-%d")
    entry = f"\n- {content}  ({timestamp})"
    p.write_text((existing + entry + "\n"), encoding="utf-8")


# ---------------- Document extraction ----------------

# Cap text payload to keep model context sane. Mirrors _read_file's 200k limit.
_EXTRACT_TEXT_CAP = 200_000

_EXTRACT_SUPPORTED = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md"}


def _parse_page_range(spec: str, total: int) -> list[int]:
    """Parse a 1-based page spec like '1-5' or '3' into 0-based indices.

    Out-of-range pages are dropped silently rather than erroring, but a
    completely empty result (or a malformed spec) raises ValueError so the
    caller sees something actionable.
    """
    s = (spec or "").strip()
    if not s:
        raise ValueError("Empty page range")
    if "-" in s:
        parts = s.split("-", 1)
        try:
            start = int(parts[0])
            end = int(parts[1])
        except ValueError as e:
            raise ValueError(f"Invalid page range '{spec}'") from e
        if start < 1 or end < start:
            raise ValueError(f"Invalid page range '{spec}'")
    else:
        try:
            start = end = int(s)
        except ValueError as e:
            raise ValueError(f"Invalid page range '{spec}'") from e
        if start < 1:
            raise ValueError(f"Invalid page range '{spec}'")
    indices = [i - 1 for i in range(start, end + 1) if 1 <= i <= total]
    if not indices:
        raise ValueError(f"Page range '{spec}' is outside the document (1..{total})")
    return indices


def _cap_text(text: str, metadata: dict[str, Any]) -> str:
    if len(text) > _EXTRACT_TEXT_CAP:
        metadata["truncated"] = True
        return text[:_EXTRACT_TEXT_CAP]
    return text


def _extract_pdf(p: Path, fmt: str, pages: str | None) -> dict[str, Any]:
    import pypdf

    reader = pypdf.PdfReader(str(p))
    total = len(reader.pages)
    indices = _parse_page_range(pages, total) if pages else list(range(total))

    metadata: dict[str, Any] = {}
    info = reader.metadata or {}
    for key in ("/Title", "/Author", "/Subject", "/Creator", "/Producer"):
        val = info.get(key)
        if val:
            metadata[key.lstrip("/").lower()] = str(val)

    chunks: list[str] = []
    for idx in indices:
        try:
            page_text = reader.pages[idx].extract_text() or ""
        except Exception:
            page_text = ""
        if page_text.strip():
            chunks.append(page_text)

    text = "\n\n".join(chunks).strip()
    if not text:
        # No extractable text typically means a scanned image PDF. Surface
        # this rather than silently returning empty so the agent can decide
        # to fall back to OCR or tell the user.
        metadata["likely_scanned"] = True

    tables: list[dict[str, Any]] = []
    # pdfplumber is heavier than pypdf, only spin it up when tables are
    # actually requested (markdown output) and the doc isn't likely scanned.
    if fmt == "markdown" and not metadata.get("likely_scanned"):
        try:
            import pdfplumber

            with pdfplumber.open(str(p)) as pdf:
                for idx in indices:
                    if idx >= len(pdf.pages):
                        continue
                    for raw in pdf.pages[idx].extract_tables() or []:
                        if not raw:
                            continue
                        tables.append({
                            "page": idx + 1,
                            "rows": [[("" if c is None else str(c)) for c in row] for row in raw],
                        })
        except Exception:
            # Table extraction is best-effort; never fail the whole call
            # because pdfplumber choked on a malformed table.
            pass

    text = _cap_text(text, metadata)
    return {
        "text": text,
        "tables": tables,
        "page_count": total,
        "metadata": metadata,
        "format": fmt,
    }


def _extract_docx(p: Path, fmt: str) -> dict[str, Any]:
    import docx

    doc = docx.Document(str(p))
    paras = [para.text for para in doc.paragraphs if para.text]
    text = "\n\n".join(paras).strip()

    metadata: dict[str, Any] = {}
    core = doc.core_properties
    if core.title:
        metadata["title"] = core.title
    if core.author:
        metadata["author"] = core.author

    text = _cap_text(text, metadata)
    return {
        "text": text,
        "tables": [],
        "page_count": 0,  # docx has no fixed page count without rendering
        "metadata": metadata,
        "format": fmt,
    }


def _extract_xlsx(p: Path, fmt: str) -> dict[str, Any]:
    import openpyxl

    wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    sheet_count = len(wb.sheetnames)
    text_chunks: list[str] = []
    tables: list[dict[str, Any]] = []
    for sheet in wb.worksheets:
        rows: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True):
            rows.append([("" if v is None else str(v)) for v in row])
        # Strip trailing all-empty rows that openpyxl can leave behind.
        while rows and not any(c.strip() for c in rows[-1]):
            rows.pop()
        if not rows:
            continue
        tables.append({"sheet": sheet.title, "rows": rows})
        text_chunks.append(f"# {sheet.title}\n" + "\n".join("\t".join(r) for r in rows))
    wb.close()

    metadata: dict[str, Any] = {"sheet_count": sheet_count}
    text = "\n\n".join(text_chunks).strip()
    text = _cap_text(text, metadata)
    return {
        "text": text,
        "tables": tables,
        "page_count": sheet_count,
        "metadata": metadata,
        "format": fmt,
    }


def _extract_pptx(p: Path, fmt: str) -> dict[str, Any]:
    import pptx

    prs = pptx.Presentation(str(p))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    slide_lines.append(line)
        if slide_lines:
            chunks.append(f"# Slide {i}\n" + "\n".join(slide_lines))

    metadata: dict[str, Any] = {}
    text = "\n\n".join(chunks).strip()
    text = _cap_text(text, metadata)
    return {
        "text": text,
        "tables": [],
        "page_count": len(prs.slides),
        "metadata": metadata,
        "format": fmt,
    }


def _extract_plaintext(p: Path, fmt: str) -> dict[str, Any]:
    raw = p.read_text(encoding="utf-8", errors="replace")
    metadata: dict[str, Any] = {}
    text = _cap_text(raw, metadata)
    return {
        "text": text,
        "tables": [],
        "page_count": 0,
        "metadata": metadata,
        "format": fmt,
    }


def _extract_document(path: str, fmt: str, pages: str | None) -> dict[str, Any]:
    p = Path(path).expanduser()
    if not p.exists():
        raise ValueError(f"File not found: {path}")
    if not p.is_file():
        raise ValueError(f"Not a file: {path}")

    ext = p.suffix.lower()
    if ext not in _EXTRACT_SUPPORTED:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(sorted(_EXTRACT_SUPPORTED))}"
        )

    if pages and ext != ".pdf":
        # Other formats don't have a stable page concept here. Better to
        # tell the caller than to silently ignore the filter.
        raise ValueError(f"'pages' is only supported for PDF, not {ext}")

    if ext == ".pdf":
        return _extract_pdf(p, fmt, pages)
    if ext == ".docx":
        return _extract_docx(p, fmt)
    if ext == ".xlsx":
        return _extract_xlsx(p, fmt)
    if ext == ".pptx":
        return _extract_pptx(p, fmt)
    return _extract_plaintext(p, fmt)


def _dctl_env() -> dict[str, str]:
    env = os.environ.copy()
    repo = Path(__file__).resolve().parents[3] / "dctl"
    if repo.exists():
        prev = env.get("PYTHONPATH", "")
        env["PYTHONPATH"] = f"{repo}{os.pathsep}{prev}" if prev else str(repo)
    return env


def _run_dctl(subcommand: str, args: list[str], cwd: str) -> dict[str, Any]:
    if not subcommand:
        raise ValueError("dctl requires a subcommand")
    cmd = [sys.executable, "-m", "dctl", subcommand, *args]
    try:
        result = subprocess.run(
            cmd,
            cwd=Path(cwd).expanduser(),
            capture_output=True,
            text=True,
            timeout=120,
            env=_dctl_env(),
        )
    except subprocess.TimeoutExpired:
        return {
            "ok": False,
            "returncode": -1,
            "output": "dctl command timed out after 120s. Try a simpler operation.",
        }
    output = result.stdout
    if result.stderr:
        output += ("\n" + result.stderr) if output else result.stderr
    if len(output) > 20_000:
        output = output[:20_000] + "\n…[truncated]"
    return {
        "ok": result.returncode == 0,
        "returncode": result.returncode,
        "output": output.strip(),
    }


# ---------------- Legacy <<TOOL>> marker parser ----------------

def parse_tool_calls(text: str) -> list[dict[str, Any]]:
    """Parse <<TOOL>>...<</TOOL>> blocks from model output (fallback)."""
    calls: list[dict[str, Any]] = []
    pattern = r"<<TOOL>>([\s\S]*?)<</TOOL>>"
    for match in re.finditer(pattern, text):
        try:
            data = json.loads(match.group(1).strip())
            if isinstance(data, dict) and "tool" in data and "params" in data:
                calls.append(data)
        except (json.JSONDecodeError, KeyError):
            continue
    return calls


# silence linters on unused import
_ = shlex
_ = os
