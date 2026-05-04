"""Tests for the extract_document agent tool."""
import tempfile
import unittest
from pathlib import Path

from sidecar.agent.tools import _extract_document


def _make_pdf(path: Path, pages: list[str]) -> None:
    """Build a tiny multi-page PDF with one short text line per page."""
    # We use pypdf's writer + reportlab-free path: pypdf can author pages
    # from a blank page + an inline content stream. To keep the test
    # dependency surface to what the project already pulls in, we build
    # each page as a minimal PDF stream and merge with pypdf.
    import pypdf
    from pypdf.generic import (
        DecodedStreamObject,
        DictionaryObject,
        NameObject,
    )

    writer = pypdf.PdfWriter()
    for text in pages:
        page = writer.add_blank_page(width=300, height=200)
        # Build a simple content stream that draws the text via a Type1
        # base font (Helvetica) — one of the 14 standard fonts that
        # readers ship by default, so no font embedding is needed.
        font = DictionaryObject({
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        })
        font_ref = writer._add_object(font)
        resources = DictionaryObject({
            NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref}),
        })
        page[NameObject("/Resources")] = resources

        stream = DecodedStreamObject()
        # Escape parens/backslashes in the literal string.
        safe = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        stream.set_data(
            f"BT /F1 12 Tf 20 100 Td ({safe}) Tj ET\n".encode("latin-1")
        )
        page[NameObject("/Contents")] = writer._add_object(stream)

    with open(path, "wb") as f:
        writer.write(f)


def _make_docx(path: Path, paragraphs: list[str]) -> None:
    import docx

    doc = docx.Document()
    for para in paragraphs:
        doc.add_paragraph(para)
    doc.save(str(path))


def _make_xlsx(path: Path, rows: list[list[object]], sheet: str = "Sheet1") -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet
    for row in rows:
        ws.append(row)
    wb.save(str(path))


def _make_pptx(path: Path, slides: list[str]) -> None:
    import pptx

    prs = pptx.Presentation()
    blank = prs.slide_layouts[6]  # blank layout
    for text in slides:
        slide = prs.slides.add_slide(blank)
        # Add a textbox at a fixed position.
        from pptx.util import Inches
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text_frame.text = text
    prs.save(str(path))


class TestExtractDocument(unittest.TestCase):
    def setUp(self) -> None:
        self._tmp = tempfile.TemporaryDirectory()
        self.tmp = Path(self._tmp.name)

    def tearDown(self) -> None:
        self._tmp.cleanup()

    # -- format coverage --------------------------------------------------

    def test_pdf_text(self) -> None:
        path = self.tmp / "doc.pdf"
        _make_pdf(path, ["Hello pdf world", "Page two body", "Page three body"])
        result = _extract_document(str(path), "markdown", None)
        self.assertEqual(result["format"], "markdown")
        self.assertEqual(result["page_count"], 3)
        self.assertIn("Hello pdf world", result["text"])
        self.assertIn("Page two body", result["text"])

    def test_pdf_pages_filter(self) -> None:
        path = self.tmp / "multi.pdf"
        _make_pdf(path, ["Alpha page", "Bravo page", "Charlie page"])
        result = _extract_document(str(path), "markdown", "2")
        self.assertEqual(result["page_count"], 3)  # total, not filtered count
        self.assertIn("Bravo page", result["text"])
        self.assertNotIn("Alpha page", result["text"])
        self.assertNotIn("Charlie page", result["text"])

    def test_pdf_pages_range(self) -> None:
        path = self.tmp / "multi.pdf"
        _make_pdf(path, ["Alpha page", "Bravo page", "Charlie page"])
        result = _extract_document(str(path), "markdown", "1-2")
        self.assertIn("Alpha page", result["text"])
        self.assertIn("Bravo page", result["text"])
        self.assertNotIn("Charlie page", result["text"])

    def test_pdf_pages_invalid(self) -> None:
        path = self.tmp / "doc.pdf"
        _make_pdf(path, ["only one"])
        with self.assertRaises(ValueError):
            _extract_document(str(path), "markdown", "5-9")
        with self.assertRaises(ValueError):
            _extract_document(str(path), "markdown", "abc")

    def test_docx_text(self) -> None:
        path = self.tmp / "doc.docx"
        _make_docx(path, ["First paragraph here", "Second paragraph here"])
        result = _extract_document(str(path), "markdown", None)
        self.assertIn("First paragraph here", result["text"])
        self.assertIn("Second paragraph here", result["text"])
        self.assertEqual(result["tables"], [])

    def test_xlsx_text_and_tables(self) -> None:
        path = self.tmp / "data.xlsx"
        _make_xlsx(path, [["name", "age"], ["Ada", 36], ["Linus", 54]])
        result = _extract_document(str(path), "markdown", None)
        self.assertIn("Ada", result["text"])
        self.assertEqual(len(result["tables"]), 1)
        self.assertEqual(result["tables"][0]["sheet"], "Sheet1")
        self.assertIn(["Ada", "36"], result["tables"][0]["rows"])

    def test_pptx_text(self) -> None:
        path = self.tmp / "deck.pptx"
        _make_pptx(path, ["Slide one body", "Slide two body"])
        result = _extract_document(str(path), "markdown", None)
        self.assertIn("Slide one body", result["text"])
        self.assertIn("Slide two body", result["text"])
        self.assertEqual(result["page_count"], 2)

    def test_txt(self) -> None:
        path = self.tmp / "notes.txt"
        path.write_text("plain text body\nsecond line", encoding="utf-8")
        result = _extract_document(str(path), "markdown", None)
        self.assertIn("plain text body", result["text"])
        self.assertIn("second line", result["text"])

    def test_md(self) -> None:
        path = self.tmp / "readme.md"
        path.write_text("# Title\n\nbody paragraph", encoding="utf-8")
        result = _extract_document(str(path), "markdown", None)
        self.assertIn("body paragraph", result["text"])

    # -- error paths ------------------------------------------------------

    def test_missing_file(self) -> None:
        with self.assertRaises(ValueError) as cm:
            _extract_document(str(self.tmp / "nope.pdf"), "markdown", None)
        self.assertIn("not found", str(cm.exception).lower())

    def test_unsupported_extension(self) -> None:
        path = self.tmp / "data.bin"
        path.write_bytes(b"\x00\x01\x02")
        with self.assertRaises(ValueError) as cm:
            _extract_document(str(path), "markdown", None)
        self.assertIn("unsupported", str(cm.exception).lower())

    def test_pages_only_for_pdf(self) -> None:
        path = self.tmp / "doc.docx"
        _make_docx(path, ["body"])
        with self.assertRaises(ValueError):
            _extract_document(str(path), "markdown", "1")


if __name__ == "__main__":
    unittest.main()
